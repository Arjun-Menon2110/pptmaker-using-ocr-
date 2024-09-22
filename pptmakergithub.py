import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import requests

# Google Custom Search API configuration
GOOGLE_API_KEY = "YOUR_GOOGLE_API_KEY"  # Replace with your Google API Key
GOOGLE_CX = "YOUR_GOOGLE_CX"  # Replace with your Google Custom Search Engine ID

def extract_text_from_image(image_path):
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return ""

def extract_questions(text):
    lines = text.split('\n')
    questions = [line.strip() for line in lines if line.strip().endswith('?')]
    return questions

def search_google(query, api_key, cx, num_results=3):
    try:
        url = "https://www.googleapis.com/customsearch/v1"
        params = {
            "key": api_key,
            "cx": cx,
            "q": query,
            "num": num_results  # Fetch more than one result for better precision
        }
        response = requests.get(url, params=params)
        response.raise_for_status()
        results = response.json()
        
        # Check for items and aggregate snippets
        if "items" in results:
            snippets = []
            for item in results["items"]:
                snippet = item.get("snippet", "")
                if snippet:
                    snippets.append(snippet)
            
            # Combine the snippets for better precision
            if snippets:
                return " ".join(snippets[:num_results])  # Combine top 'num_results' snippets
            else:
                return "No detailed answer found."
        else:
            return "No answer found."
    except Exception as e:
        print(f"Error searching Google: {e}")
        return "Error getting answer."

def create_ppt(question_answer_pairs, output_path="formatted_presentation.pptx"):
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    font_size_title = Pt(32)
    font_size_text = Pt(20)
    
    for question, answer in question_answer_pairs:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Format title (bold and centered)
        title = slide.shapes.title
        title.text = "Question and Answer"
        title.text_frame.paragraphs[0].font.size = font_size_title
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add question box with space between question and answer
        question_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), slide_width - Inches(2), Inches(2))
        question_frame = question_box.text_frame
        question_frame.text = "Question: " + question
        question_frame.paragraphs[0].font.size = font_size_text
        question_frame.paragraphs[0].font.bold = True
        
        # Add answer box with more space and better formatting
        answer_box = slide.shapes.add_textbox(Inches(1), Inches(4), slide_width - Inches(2), Inches(2.5))
        answer_frame = answer_box.text_frame
        answer_frame.text = "Answer: " + answer
        answer_frame.paragraphs[0].font.size = font_size_text
        answer_frame.paragraphs[0].font.bold = False
    
    prs.save(output_path)
    print(f"PowerPoint presentation saved at: {output_path}")

def main(image_path, api_key, cx):
    text = extract_text_from_image(image_path)
    print(f"Extracted Text: {text}")
    
    questions = extract_questions(text)
    print(f"Extracted Questions: {questions}")
    
    question_answer_pairs = []
    for q in questions:
        print(f"Processing question: {q}")
        answer = search_google(q, api_key, cx)
        print(f"Answer: {answer}")
        question_answer_pairs.append((q, answer))
    
    create_ppt(question_answer_pairs, "questions_answers_formatted_presentation.pptx")

# Run the main function with your image file and API credentials
image_path =   # Path to your uploaded image
api_key =   # Your Google API Key
cx =  # Your Google Custom Search Engine ID
main(image_path, api_key, cx)